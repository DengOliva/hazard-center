"""Generate a filled PPTX from the template + database data + manual entries."""
import json
import sqlite3
from copy import deepcopy
from datetime import date, datetime, timedelta
from io import BytesIO
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parent
DATA_DIR = Path(__import__("os").environ.get("DATA_DIR", ROOT / "data"))
DB_PATH = DATA_DIR / "hazards.db"
TEMPLATE_PATH = ROOT / "体系部门会议材料7.2.pptx"
MEETING_DATA_FILE = DATA_DIR / "meeting_data.json"


def _db():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn


def _tbl_rows(table):
    """Return list of lists for a table (excluding header row 0)."""
    rows = []
    for r in range(1, len(table.rows)):
        row = [table.cell(r, c).text.strip() for c in range(len(table.columns))]
        rows.append(row)
    return rows


def _set_cell(table, row, col, text):
    """Set cell text, preserving first paragraph formatting."""
    cell = table.cell(row, col)
    p = cell.text_frame.paragraphs[0]
    p.text = str(text)
    # Clear any extra paragraphs
    for extra in cell.text_frame.paragraphs[1:]:
        extra.text = ""


def _insert_row(table, after_row, n=1):
    """Insert n empty rows after the given row index."""
    for _ in range(n):
        # Copy the last row's XML
        tr = table._tbl
        new_tr = deepcopy(tr[-1])
        # Insert before the last row (which we'll keep as template)
        tr.insert(len(tr) - 1, new_tr)


def _remove_row(table, row_idx):
    """Remove a row from the table."""
    tr = table._tbl
    if 0 <= row_idx < len(tr):
        tr.remove(tr[row_idx])


def _update_title(slide, text, title_shape_name=None):
    """Update the title text on a content slide."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            full = tf.text.strip()
            # Look for shape containing section number + title pattern
            if full and len(full) > 5:
                # Replace the date portion: everything after （ or （
                for para in tf.paragraphs:
                    for run in para.runs:
                        if "—" in run.text or "--" in run.text:
                            # Replace date range part
                            idx = run.text.find("（")
                            if idx >= 0:
                                run.text = run.text[:idx] + text
                            else:
                                run.text = text
                            return
            if title_shape_name and shape.name == title_shape_name:
                tf.paragraphs[0].text = text
                return


def _fill_auto_slides(prs, data, week_range):
    """Fill slides with auto-generated data."""
    slides = list(prs.slides)

    # ─── Slide 6: 监查/监督报告 ─── (manual table, fill from saved)
    _fill_manual_table(slides[5], 'slide_6', 0, 1)

    # ─── Slide 7: 行为偏差统计 ───
    _fill_slide7(slides[6], data)

    # ─── Slide 8: 隐患记录提交 ───
    _fill_slide8(slides[7], data)

    # ─── Slide 11: 预警刹车 ───
    _fill_slide11(slides[10], data)

    # ─── Slides 12-14: 整改单/通报/停工令关闭 ───
    _fill_close_slide(slides[11], data, 'rectification')
    _fill_close_slide(slides[12], data, 'notice')
    _fill_close_slide(slides[13], data, 'stop_work')

    # ─── Slide 16: 培训 ───
    _fill_slide16(slides[15], data, week_range)

    # ─── Manual slides ───
    _fill_manual_table(slides[2], 'slide_3', 0, 1)   # 年度计划6月
    _fill_manual_table(slides[3], 'slide_4', 0, 1)   # 上级关注
    _fill_manual_table(slides[4], 'slide_5', 0, 1)   # 年度计划7月
    _fill_manual_table(slides[8], 'slide_9', 0, 1)   # 作业风险管控
    _fill_manual_table(slides[9], 'slide_10', 0, 1)  # 隐患抽查
    _fill_manual_table(slides[14], 'slide_15', 0, 1) # 证件管理
    _fill_manual_table(slides[16], 'slide_17', 0, 1) # 体系工作一
    _fill_manual_table(slides[17], 'slide_18', 0, 1) # 体系工作二


def _fill_slide7(slide, data):
    """Slide 7: 内外部行为偏差统计 — 2 tables."""
    bs = data.get("behavior_stats", {})
    internal = bs.get("internal", [])
    external = bs.get("external", [])
    tables = [s for s in slide.shapes if s.has_table]
    if len(tables) < 2:
        return
    # Left table = external (工程公司), Right table = internal (内部)
    # Actually template has left=工程公司, right=内部录入
    # Let's find by header text
    for tbl_shape in tables:
        tbl = tbl_shape.table
        hdr = tbl.cell(0, 0).text.strip()
        if "工程公司" in hdr:
            items = external
        elif "内部" in hdr:
            items = internal
        else:
            continue

        # Ensure enough data rows (keep header row 1 as template)
        data_start_row = 2  # row 0=title, row 1=headers, row 2+=data
        needed = len(items) + 1  # +1 for 合计 row
        current_data_rows = len(tbl.rows) - data_start_row

        if needed > current_data_rows:
            for _ in range(needed - current_data_rows):
                _insert_row(tbl, len(tbl.rows) - 1)
        elif needed < current_data_rows:
            for _ in range(current_data_rows - needed):
                _remove_row(tbl, data_start_row + needed)

        for i, item in enumerate(items):
            r = data_start_row + i
            if r < len(tbl.rows):
                _set_cell(tbl, r, 0, str(i + 1))
                _set_cell(tbl, r, 1, item["type"])
                _set_cell(tbl, r, 2, str(item["count"]))
                _set_cell(tbl, r, 3, "")
        # 合计 row
        total_row = data_start_row + len(items)
        if total_row < len(tbl.rows):
            _set_cell(tbl, total_row, 1, "合计")
            _set_cell(tbl, total_row, 2, str(sum(it["count"] for it in items)))

    # Update summary text
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name.startswith("文本"):
            tf = shape.text_frame
            int_sum = sum(it["count"] for it in internal)
            ext_sum = sum(it["count"] for it in external)
            tf.paragraphs[0].text = (
                f"本期内部行为偏差共{int_sum}项，"
                f"工程公司录入行为偏差共{ext_sum}项。"
            )


def _fill_slide8(slide, data):
    """Slide 8: 隐患记录提交 — main table + contractor table + total."""
    ps = data.get("people_stats", {})
    groups = ps.get("groups", [])
    if not groups:
        return

    tables = [s for s in slide.shapes if s.has_table]
    if not tables:
        return

    # Find the main table (the largest one, 7 columns)
    main_tbl = None
    sub_tbl = None
    total_tbl = None
    for s in tables:
        tbl = s.table
        nc = len(tbl.columns)
        if nc == 7:
            main_tbl = tbl
        elif nc == 6:
            sub_tbl = tbl
        elif nc == 4:
            total_tbl = tbl

    overall_total = ps.get("total_overall", 0)

    if total_tbl and len(total_tbl.rows) >= 1:
        _set_cell(total_tbl, 0, 3, str(overall_total))

    if main_tbl:
        data_start = 2  # row 0=merged title, row 1=headers
        # Build rows for each person
        rows_data = []
        for g in groups:
            for m in g.get("members", []):
                rows_data.append({
                    "dept": g["category"],
                    "name": m["name"],
                    "count": str(m["count"]),
                    "avg": f"{m['count'] / 7:.1f}",
                    "met": "是" if m["count"] / 7 >= 1 else "否",
                })
            rows_data.append({"dept": f"{g['category']} 合计", "name": "", "count": str(g["total"]), "avg": "", "met": ""})

        needed = len(rows_data)
        current = len(main_tbl.rows) - data_start
        if needed > current:
            # Insert more rows
            for _ in range(needed - current):
                _insert_row(main_tbl, len(main_tbl.rows) - 1)
        elif needed < current:
            for _ in range(current - needed):
                _remove_row(main_tbl, data_start + needed)

        for i, rd in enumerate(rows_data):
            r = data_start + i
            if r >= len(main_tbl.rows):
                break
            _set_cell(main_tbl, r, 0, rd["dept"])
            _set_cell(main_tbl, r, 1, rd["name"])
            _set_cell(main_tbl, r, 2, rd["count"])
            _set_cell(main_tbl, r, 3, rd["avg"])
            _set_cell(main_tbl, r, 4, rd["met"])
            _set_cell(main_tbl, r, 5, "")


def _fill_slide11(slide, data):
    """Slide 11: 管理干部预警刹车 — 13-column table."""
    da = data.get("dept_alert", {})
    depts = da.get("depts", [])
    int_types = da.get("int_types", [])
    ext_types = da.get("ext_types", [])

    tables = [s for s in slide.shapes if s.has_table]
    if not tables or not depts:
        return
    tbl = tables[0].table

    data_start = 3  # row 0=merged title, row 1=sub-header, row 2=stat headers
    needed = len(depts) + 1  # +1 for total row
    current = len(tbl.rows) - data_start

    if needed > current:
        for _ in range(needed - current):
            _insert_row(tbl, len(tbl.rows) - 1)
    elif needed < current:
        for _ in range(current - needed):
            _remove_row(tbl, data_start + needed)

    for i, d in enumerate(depts):
        r = data_start + i
        if r >= len(tbl.rows):
            break
        _set_cell(tbl, r, 0, d.get("dept", ""))
        for j, t in enumerate(int_types):
            _set_cell(tbl, r, 1 + j, str(d.get(f"int_{t}", 0) or 0))
        offset = 1 + len(int_types)
        for j, t in enumerate(ext_types):
            _set_cell(tbl, r, offset + j, str(d.get(f"ext_{t}", 0) or 0))


def _fill_close_slide(slide, data, key):
    """Slide 12/13/14: Close status tables."""
    cs = data.get("close_status", {})
    items = cs.get(key, [])
    if not items:
        return

    tables = [s for s in slide.shapes if s.has_table]
    if not tables:
        return
    tbl = tables[0].table

    data_start = 1  # row 0=headers
    needed = len(items)
    current = len(tbl.rows) - data_start

    if needed > current:
        for _ in range(needed - current):
            _insert_row(tbl, len(tbl.rows) - 1)
    elif needed < current:
        for _ in range(current - needed):
            _remove_row(tbl, data_start + needed)

    cols = len(tbl.columns)
    for i, item in enumerate(items):
        r = data_start + i
        if r >= len(tbl.rows):
            break
        _set_cell(tbl, r, 0, str(i + 1))
        if cols > 1:
            _set_cell(tbl, r, 1, item.get("source", ""))
        if cols > 2:
            _set_cell(tbl, r, 2, item.get("dept", ""))
        if cols > 3:
            _set_cell(tbl, r, 3, str(item.get("no", item.get("content", "")))[:40])
        if cols > 4:
            _set_cell(tbl, r, 4, item.get("issue_date", ""))
        # Use remaining columns for content/deadline/closed/remark as available
        closed = "是" if item.get("closed") in ("是", "已关闭") else "否"

    # Update summary text
    closed_count = sum(1 for it in items if it.get("closed") in ("是", "已关闭"))
    open_count = len(items) - closed_count
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name.startswith("文本"):
            name_map = {"rectification": "整改通知单", "notice": "通报", "stop_work": "停工令"}
            label = name_map.get(key, key)
            shape.text_frame.paragraphs[0].text = (
                f"本期共收到{label}{len(items)}份，已关闭{closed_count}份，未关闭{open_count}份。"
            )


def _fill_slide16(slide, data, week_range):
    """Slide 16: 培训."""
    training = data.get("training", {})
    items = training.get("items", [])
    if not items:
        return

    tables = [s for s in slide.shapes if s.has_table]
    if not tables:
        return
    tbl = tables[0].table

    data_start = 1
    needed = len(items)
    current = len(tbl.rows) - data_start

    if needed > current:
        for _ in range(needed - current):
            _insert_row(tbl, len(tbl.rows) - 1)
    elif needed < current:
        for _ in range(current - needed):
            _remove_row(tbl, data_start + needed)

    for i, item in enumerate(items):
        r = data_start + i
        if r >= len(tbl.rows):
            break
        _set_cell(tbl, r, 0, str(i + 1))
        _set_cell(tbl, r, 1, item.get("date", ""))
        _set_cell(tbl, r, 2, item.get("content", ""))
        _set_cell(tbl, r, 3, "")


def _fill_manual_table(slide, key, data_start_row, header_rows):
    """Fill a manual-entry table from saved meeting data."""
    try:
        manual = json.loads(MEETING_DATA_FILE.read_text(encoding="utf-8"))
    except Exception:
        return
    saved = manual.get(key, {})
    table_data = saved.get("table", [])
    if not table_data:
        return

    tables = [s for s in slide.shapes if s.has_table]
    if not tables:
        return
    tbl = tables[0].table

    data_start = data_start_row + header_rows
    needed = len(table_data)
    current = len(tbl.rows) - data_start

    if needed > current:
        for _ in range(needed - current):
            _insert_row(tbl, len(tbl.rows) - 1)
    elif needed < current:
        for _ in range(current - needed):
            _remove_row(tbl, data_start + needed)

    for ri, row in enumerate(table_data):
        r = data_start + ri
        if r >= len(tbl.rows):
            break
        for ci, cell_val in enumerate(row):
            if ci < len(tbl.columns):
                _set_cell(tbl, r, ci, str(cell_val or ""))


def generate(start_date, end_date):
    """Generate a PPTX file from template with all data filled in.

    Returns BytesIO buffer.
    """
    if not TEMPLATE_PATH.exists():
        raise FileNotFoundError(f"Template not found: {TEMPLATE_PATH}")

    prs = Presentation(str(TEMPLATE_PATH))

    # Load auto data (same logic as meeting.py's /api/meeting/data)
    from meeting import _load_alert_data, _build_close_status, _build_training_summary

    alert_records = _load_alert_data()

    # Behavior stats
    behavior_stats = {"internal": [], "external": []}
    if alert_records:
        type_int, type_ext = {}, {}
        for rec in alert_records:
            if rec["date"] < start_date or rec["date"] > end_date:
                continue
            tn = rec["type_name"]
            if rec["category"] == "internal":
                type_int[tn] = type_int.get(tn, 0) + 1
            else:
                type_ext[tn] = type_ext.get(tn, 0) + 1
        behavior_stats["internal"] = [{"type": k, "count": v} for k, v in type_int.items()]
        behavior_stats["external"] = [{"type": k, "count": v} for k, v in type_ext.items()]

    # People stats
    with _db() as conn:
        people_rows = conn.execute(
            "SELECT * FROM people WHERE active=1 ORDER BY category, department, name"
        ).fetchall()
        counts = {
            r[0]: r[1] for r in conn.execute(
                "SELECT checker_name, COUNT(*) FROM hazards WHERE check_date BETWEEN ? AND ? GROUP BY checker_name",
                (start_date, end_date),
            )
        }
    groups = {}
    for p in people_rows:
        count = counts.get(p["name"], 0)
        key = p["category"]
        if key not in groups:
            groups[key] = {"category": key, "members": [], "total": 0}
        groups[key]["members"].append({"name": p["name"], "category": key, "department": p["department"], "count": count})
        groups[key]["total"] += count

    people_stats = {"groups": list(groups.values()), "total_overall": sum(g["total"] for g in groups.values())}

    # Dept alert stats
    dept_alert = {"ext_types": [], "int_types": [], "depts": []}
    if alert_records:
        ext_types = ["挂牌督办", "管理约谈", "红黄牌", "处理通报", "停工令", "整改单", "违章培训通知单", "监理通知单"]
        int_types = ["停工令", "处理通报", "整改单", "违章培训通知单"]
        dept_map = {}
        for rec in alert_records:
            if rec["date"] < start_date or rec["date"] > end_date:
                continue
            if not rec["dept_name"]:
                continue
            dept = rec["dept_name"]
            if dept not in dept_map:
                dept_map[dept] = {"dept": dept}
                for t in ext_types:
                    dept_map[dept][f"ext_{t}"] = 0
                for t in int_types:
                    dept_map[dept][f"int_{t}"] = 0
            key = f"ext_{rec['type_name']}" if rec["category"] == "external" else f"int_{rec['type_name']}"
            if key in dept_map[dept]:
                dept_map[dept][key] += 1
        dept_alert = {"ext_types": ext_types, "int_types": int_types, "depts": list(dept_map.values())}

    close_status = _build_close_status(alert_records, start_date, end_date)
    training = _build_training_summary(start_date, end_date)

    data = {
        "behavior_stats": behavior_stats,
        "people_stats": people_stats,
        "dept_alert": dept_alert,
        "close_status": close_status,
        "training": training,
    }

    week_range = f"（{start_date}--{end_date}）"
    _fill_auto_slides(prs, data, week_range)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
